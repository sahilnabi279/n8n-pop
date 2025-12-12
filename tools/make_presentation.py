from pptx import Presentation
import json
prs = Presentation()
# Title slide
s = prs.slides.add_slide(prs.slide_layouts[0])
s.shapes.title.text = "n8n Workflow Popularity — Summary"
s.placeholders[1].text = "YouTube + Forum data merged"

# Methodology
s = prs.slides.add_slide(prs.slide_layouts[1])
s.shapes.title.text = "Methodology"
s.placeholders[1].text = "Collectors: YouTube API, Discourse forum. Merge by normalized titles. Summary metrics: total_views, avg_like_to_view_ratio."

# Data sources
s = prs.slides.add_slide(prs.slide_layouts[1])
s.shapes.title.text = "Data Sources"
s.placeholders[1].text = "YouTube (views/likes/comments)\\nDiscourse (replies/views/contributors)"

# Data counts
s = prs.slides.add_slide(prs.slide_layouts[1])
s.shapes.title.text = "Dataset size"
with open('data/final_workflows.json','r',encoding='utf-8') as f:
    items = json.load(f)
s.placeholders[1].text = f"Final merged workflows: {len(items)}"

# Top 5 by views
s = prs.slides.add_slide(prs.slide_layouts[1])
s.shapes.title.text = "Top 5 by total_views"
tops = sorted(items, key=lambda x: x.get('summary',{}).get('total_views',0), reverse=True)[:5]
body = "\\n".join([f\"{i+1}. {t['workflow']} — {t['summary'].get('total_views',0)} views\" for i,t in enumerate(tops)])
s.placeholders[1].text = body

# Next steps
s = prs.slides.add_slide(prs.slide_layouts[1])
s.shapes.title.text = "Next Steps"
s.placeholders[1].text = "Add more sources; retry Trends with smaller batches; deploy API."

prs.save('n8n_popularity_presentation.pptx')
print('Saved n8n_popularity_presentation.pptx')
